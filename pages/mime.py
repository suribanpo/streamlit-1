import streamlit as st
from pptx import Presentation
import io
from fpdf import FPDF
import pandas as pd
from docx import Document

# 제목 설정
st.title("What is Mime?")

# 사용자 입력 받기
slide_title = st.text_input("강의 제목", "강의 제목을 넣어주세요.")
slide_subtitle = st.text_input("부제목", "강의자의 이름 등을 넣어주세요.")
class_outline = st.text_area("작성할 내용을 적어주세요.")

# 파워포인트 파일 생성
presentation = Presentation()
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.placeholders[0]
title.text = slide_title
subtitle = slide.placeholders[1]
subtitle.text = slide_subtitle
lines = class_outline.strip().split("\n")
for line in lines:
    line = line.strip()
    if line:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = line

ppt = io.BytesIO()
presentation.save(ppt)
ppt.seek(0)
# CSV 파일 생성
df = pd.DataFrame({
    "제목": [slide_title],
    "부제목": [slide_subtitle],
    "내용": [class_outline]
})

# StringIO를 사용하여 CSV 문자열로 저장
csv_output = io.StringIO()
df.to_csv(csv_output, index=False)
csv_output.seek(0)  # 스트림을 시작 위치로 이동

# 다운로드 버튼 생성
st.download_button("CSV 다운로드", data=csv_output.getvalue(), file_name="강의.csv", mime='text/csv')

# CSV 파일 생성
df = pd.DataFrame({
    "제목": [slide_title],
    "부제목": [slide_subtitle],
    "내용": [class_outline]
})
csv_output = io.StringIO()
df.to_csv(csv_output, index=False)
csv_output.seek(0)

# 워드 파일 생성
doc = Document()
doc.add_heading(slide_title, level=0)
doc.add_heading(slide_subtitle, level=1)
doc.add_paragraph(class_outline)
word = io.BytesIO()
doc.save(word)
word.seek(0)

# 다운로드 버튼들
st.download_button("파워포인트 다운로드", data=ppt, file_name="강의.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
