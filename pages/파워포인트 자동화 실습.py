import streamlit as st
from pptx import Presentation

st.title("파워포인트 작성하기")
# 강의 제목 및 부제목
slide_title = "이차방정식의 활용"
slide_subtitle = '황수빈 선생님의 수학 여행'

# 강의 개요 확장
class_outline = st.text_area("개요를 작성해주세요.")

# 출력하여 확인
st.write(f"강의 제목: {slide_title}")
st.write(f"부제목: {slide_subtitle}")
st.write("강의 개요:")
st.write(class_outline)



# Create a PowerPoint presentation
presentation = Presentation()

title_slide_layout = presentation.slide_layouts[0] # 0 : 제목슬라이드에 해당
slide = presentation.slides.add_slide(title_slide_layout) # 슬라이드 추가

# 제목 - 제목에 값넣기
title = slide.placeholders[0] # 제목
title.text = slide_title # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목상자는 placeholders[0], 부제목상자는 [1]
subtitle.text = slide_subtitle

# Split the class outline into lines
lines = class_outline.strip().split("\n")

# Iterate through the remaining lines to create slides
for i in range(len(lines)):
    line = lines[i].strip()
    if line.startswith("-"):
        # Add the line to the current slide's body
        body_shape = slide.shapes.placeholders[1]
        p = body_shape.text_frame.add_paragraph()
        p.text = line.strip("- ").strip()
    else:
        # Create a new slide with the topic as the title
        slide_layout = presentation.slide_layouts[1]  # Content slide layout
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = line.strip()

# presentation.save("이차방정식 자료.pptx")

import io

ppt = io.BytesIO()
presentation.save(ppt)
ppt.seek(0)
st.download_button(label = "다운로드하기", data = ppt,  file_name = "개요.pptx", mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation")
